from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_test_presentation():
    """Create a test PowerPoint presentation"""
    print("Creating test PowerPoint presentation...")
    
    # Create presentation
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Smart Healthcare Mobile App"
    subtitle.text = "AI-Powered Rural Healthcare Solution\nTeam: HealthTech Innovators\nSIH 2024"
    
    # Slide 2: Problem Statement
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Problem Statement"
    tf = content.text_frame
    tf.text = "Rural areas face significant healthcare challenges:"
    
    p = tf.add_paragraph()
    p.text = "• Limited access to healthcare professionals"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Inadequate medical record management"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Lack of telemedicine infrastructure"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Poor connectivity and digital literacy"
    p.level = 1
    
    # Slide 3: Our Solution
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Our Solution: HealthConnect App"
    tf = content.text_frame
    tf.text = "Comprehensive mobile healthcare platform:"
    
    features = [
        "Patient registration and profile management",
        "Appointment booking with healthcare providers", 
        "Secure medical record storage and sharing",
        "Video consultation and telemedicine",
        "Medicine reminders and health tracking",
        "Offline functionality for low connectivity areas"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = f"• {feature}"
        p.level = 1
    
    # Slide 4: Technology Stack
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technology Stack"
    tf = content.text_frame
    tf.text = "Modern, scalable architecture:"
    
    tech_stack = [
        "Frontend: React Native (cross-platform)",
        "Backend: Node.js with Express framework", 
        "Database: MongoDB for flexible data storage",
        "Cloud: AWS for infrastructure and scalability",
        "Authentication: JWT tokens and OAuth 2.0",
        "Video: WebRTC for real-time communication"
    ]
    
    for tech in tech_stack:
        p = tf.add_paragraph()
        p.text = f"• {tech}"
        p.level = 1
    
    # Slide 5: System Architecture
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "System Architecture"
    tf = content.text_frame
    tf.text = "Microservices-based architecture ensuring scalability and maintainability:"
    
    architecture = [
        "API Gateway for request routing and authentication",
        "User Service for patient and provider management",
        "Appointment Service for scheduling and notifications",
        "Medical Records Service with encryption",
        "Video Service for telemedicine consultations",
        "Notification Service for reminders and alerts"
    ]
    
    for arch in architecture:
        p = tf.add_paragraph()
        p.text = f"• {arch}"
        p.level = 1
    
    # Slide 6: Implementation Timeline
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Implementation Timeline"
    tf = content.text_frame
    tf.text = "4-month development roadmap:"
    
    timeline = [
        "Phase 1 (Month 1-2): Core app development and user management",
        "Phase 2 (Month 3): Telemedicine integration and video calling",
        "Phase 3 (Month 4): Testing, optimization, and deployment",
        "Post-launch: Continuous monitoring and feature updates"
    ]
    
    for phase in timeline:
        p = tf.add_paragraph()
        p.text = f"• {phase}"
        p.level = 1
    
    # Slide 7: Team & Budget
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Team Structure & Budget"
    tf = content.text_frame
    tf.text = "Experienced team with realistic budget:"
    
    team_budget = [
        "Team: 6 developers, 1 designer, 1 project manager",
        "Development cost: $50,000 for 4 months",
        "Infrastructure: $500/month (AWS, third-party APIs)",
        "Total budget: $55,000 initial + $500/month operational"
    ]
    
    for item in team_budget:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 1
    
    # Slide 8: Supporting Links
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Supporting Materials"
    tf = content.text_frame
    tf.text = "Project resources and demonstrations:"
    
    links = [
        "GitHub Repository: https://github.com/healthtech/healthcare-app",
        "Demo Video: https://youtube.com/watch?v=healthapp-demo",
        "Technical Documentation: https://docs.google.com/document/healthcare-specs",
        "Live Prototype: https://healthconnect-demo.netlify.app"
    ]
    
    for link in links:
        p = tf.add_paragraph()
        p.text = f"• {link}"
        p.level = 1
    
    # Save presentation
    filename = "test_healthcare_presentation.pptx"
    prs.save(filename)
    print(f"Test presentation created: {filename}")
    return filename

if __name__ == "__main__":
    create_test_presentation()